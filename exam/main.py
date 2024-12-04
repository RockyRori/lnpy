import os
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document


# 提取PDF文件中的文本
def extract_text_from_pdf(pdf_path):
    try:
        with open(pdf_path, "rb") as file:
            reader = PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text()
            return text
    except Exception as e:
        return f"Error extracting text from PDF: {e}"


# 提取PPT文件中的文本
def extract_text_from_ppt(ppt_path):
    try:
        prs = Presentation(ppt_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error extracting text from PPT: {e}"


# 提取Word文件中的文本
def extract_text_from_word(docx_path):
    try:
        doc = Document(docx_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        return f"Error extracting text from Word: {e}"


# 分割并保存为多个txt文件
def save_text_in_chunks(text, base_output_path):
    max_chunk_size = 5000  # 每个txt文件的最大字数
    chunks = [text[i:i + max_chunk_size] for i in range(0, len(text), max_chunk_size)]

    for idx, chunk in enumerate(chunks):
        output_path = f"{base_output_path}_part{idx + 1}.txt"
        try:
            with open(output_path, "w", encoding="utf-8") as txt_file:
                txt_file.write(chunk)
            print(f"Text saved to {output_path}")
        except Exception as e:
            print(f"Error saving text to {output_path}: {e}")


# 提取文件并分割保存为多个txt文件
def extract_and_save_text(file_path, base_output_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif ext == ".pptx":
        text = extract_text_from_ppt(file_path)
    elif ext == ".docx":
        text = extract_text_from_word(file_path)
    else:
        return f"Unsupported file format: {ext}"

    save_text_in_chunks(text, base_output_path)


# 示例用法
file_path = "Database Management Lecture 1.pdf"  # 输入文件路径
base_output_path = "output"  # 输出文件基础路径
extract_and_save_text(file_path, base_output_path)
